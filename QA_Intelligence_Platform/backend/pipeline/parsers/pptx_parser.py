"""PowerPoint parser — extracts text from each slide."""
from __future__ import annotations
import os


def parse_pptx(path: str) -> tuple[str, dict]:
    """
    Extract all text from a .pptx file, one block per slide.
    Returns (full_text, metadata).
    """
    from pptx import Presentation

    prs = Presentation(path)
    slides_text: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = " ".join(run.text for run in para.runs if run.text.strip())
                if line:
                    parts.append(line)
        if parts:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(parts))

    full_text = "\n\n".join(slides_text)
    metadata = {
        "source":     "pptx",
        "filename":   os.path.basename(path),
        "path":       path,
        "slide_count": len(prs.slides),
    }
    return full_text, metadata
